# VERSION: 1.0.2 (Fixed Float Inches)
import os
from pptx import Presentation

from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE
from pptx.oxml.xmlchemy import OxmlElement

# Rich text support
from .rich_text_utils import apply_rich_text_formatting, has_rich_text

class QuizRenderer:
    def __init__(self, output_file='quiz_presentation.pptx', subject="通用"):
        self.output_file = output_file
        self.prs = Presentation()
        self.subject = subject
        
        # --- Default Design Tokens ---
        self.BG_COLOR = RGBColor(241, 245, 249) # #f1f5f9
        self.CARD_BG_COLOR = RGBColor(255, 255, 255) # #ffffff
        self.ACCENT_COLOR = RGBColor(59, 130, 246) # #3b82f6 (Blue)
        self.ACCENT_DARK = RGBColor(30, 64, 175) # #1e40af (Dark Blue)
        self.TEXT_STEM = RGBColor(15, 23, 42) # #0f172a (Slate 900)
        self.TEXT_LIGHT = RGBColor(51, 65, 85) # #334155 (Slate 700)
        self.TEXT_HINT = RGBColor(148, 163, 184) # #94a3b8 (Slate 400)
        self.RED_ANSWER = RGBColor(239, 68, 68) # #ef4444
        self.ANALYSIS_BG = RGBColor(239, 246, 255) # #eff6ff
        
        # --- Subject-Specific Theme Overrides ---
        self._apply_theme()
        
        self.FONT_MAIN = "Microsoft YaHei" 
        
        # Dimensions
        self.SLIDE_WIDTH = Inches(13.333)
        self.SLIDE_HEIGHT = Inches(7.5)
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT

    def _apply_theme(self):
        """Sets color tokens based on the subject."""
        themes = {
            "语文": {"bg": (254, 250, 240), "accent": (139, 69, 19), "dark": (92, 51, 23)}, # Warm Brown/Sepia
            "历史": {"bg": (254, 250, 240), "accent": (139, 69, 19), "dark": (92, 51, 23)},
            "数学": {"bg": (240, 249, 255), "accent": (37, 99, 235), "dark": (30, 58, 138)}, # Tech Blue
            "物理": {"bg": (240, 249, 255), "accent": (37, 99, 235), "dark": (30, 58, 138)},
            "化学": {"bg": (240, 249, 255), "accent": (37, 99, 235), "dark": (30, 58, 138)},
            "生物": {"bg": (240, 253, 244), "accent": (22, 163, 74), "dark": (20, 83, 45)}, # Nature Green
            "地理": {"bg": (240, 253, 244), "accent": (22, 163, 74), "dark": (20, 83, 45)},
            "英语": {"bg": (245, 243, 255), "accent": (124, 58, 237), "dark": (76, 29, 149)}, # Modern Purple
            "政治": {"bg": (254, 242, 242), "accent": (220, 38, 38), "dark": (127, 29, 29)}, # Logic Red
        }
        
        theme = themes.get(self.subject, {})
        if theme:
            self.BG_COLOR = RGBColor(*theme["bg"])
            self.ACCENT_COLOR = RGBColor(*theme["accent"])
            self.ACCENT_DARK = RGBColor(*theme["dark"])
            # Use background for analysis bg but slightly modified
            self.ANALYSIS_BG = RGBColor(
                max(0, theme["bg"][0] - 10),
                max(0, theme["bg"][1] - 10),
                max(0, theme["bg"][2] - 10)
            )

    def _set_bg(self, slide):
        """Sets the slide background color and adds decorative 'Canvas-style' elements."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.BG_COLOR
        
        # Add decorative shapes based on subject (simulating Canvas)
        if self.subject in ["数学", "物理", "化学"]:
            # Tech Grid / Geometry
            for i in range(15):
                line_w = self.SLIDE_WIDTH / 15
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 
                    i * line_w, 0, Inches(0.01), self.SLIDE_HEIGHT
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.ACCENT_COLOR
                shape.fill.fore_color.brightness = 0.85
                shape.line.visible = False
        
        elif self.subject in ["生物", "地理"]:
            # Organic Circles
            import random
            random.seed(42) # Deterministic for this slide
            for _ in range(8):
                size = Inches(random.uniform(1.5, 4))
                # Generate x, y carefully
                x_pos = random.uniform(0, 13) * 914400
                y_pos = random.uniform(0, 7) * 914400
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    int(x_pos), int(y_pos), size, size
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.ACCENT_COLOR
                shape.fill.fore_color.brightness = 0.92
                shape.line.visible = False

        elif self.subject in ["语文", "历史"]:
            # Enhanced Vintage/Humanities Border (Canva-inspired)
            margin = Inches(0.2)
            
            # Outer decorative border
            outer = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                margin, margin, 
                self.SLIDE_WIDTH - 2*margin, 
                self.SLIDE_HEIGHT - 2*margin
            )
            outer.fill.background()
            outer.line.color.rgb = self.ACCENT_DARK
            outer.line.width = Pt(3)
            
            # Inner accent line
            inner_margin = Inches(0.25)
            inner = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                inner_margin, inner_margin,
                self.SLIDE_WIDTH - 2*inner_margin,
                self.SLIDE_HEIGHT - 2*inner_margin
            )
            inner.fill.background()
            inner.line.color.rgb = self.ACCENT_COLOR
            inner.line.width = Pt(1)
            
            # Corner decorations (small accent rectangles)
            corner_size = Inches(0.15)
            corner_positions = [
                (margin, margin),  # Top-left
                (self.SLIDE_WIDTH - margin - corner_size, margin),  # Top-right
                (margin, self.SLIDE_HEIGHT - margin - corner_size),  # Bottom-left
                (self.SLIDE_WIDTH - margin - corner_size, self.SLIDE_HEIGHT - margin - corner_size)  # Bottom-right
            ]
            for x, y in corner_positions:
                corner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, corner_size, corner_size)
                corner.fill.solid()
                corner.fill.fore_color.rgb = self.ACCENT_DARK
                corner.line.visible = False

    def _add_card_container(self, slide):
        """Adds the white card with shadow and top bar."""
        # 1. Main Card
        # Approx 1280x720 in CSS inside a larger margin. 
        # In PPT, let's make it cover most of the slide with margins.
        margin_x = Inches(0.5)
        margin_y = Inches(0.6)
        width = self.SLIDE_WIDTH - (margin_x * 2)
        height = self.SLIDE_HEIGHT - (margin_y * 2)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            margin_x, margin_y, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.CARD_BG_COLOR
        shape.line.fill.background() # No border
        shape.shadow.inherit = False # Basic shadow isn't great in python-pptx, but let's leave default or none
        
        # 2. Top Gradient Bar (Simulated with Solid Blue for now or shape)
        # Position: Absolute top of card
        bar_height = Inches(0.15)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            margin_x, margin_y, width, bar_height
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.ACCENT_COLOR
        bar.line.fill.background()
        
        # 3. Bottom Right Accent Circle (Decorative)
        circle_size = Inches(3)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            self.SLIDE_WIDTH + Inches(0.5), self.SLIDE_HEIGHT + Inches(0.5), # Off canvas? No, clip
            circle_size, circle_size
        )
        # Actually place it at bottom right of SLIDE, partially visible
        # PPT doesn't clip "overflow:hidden" like CSS.
        # So we place it carefully or skip if it looks messy.
        # Let's place it inside the card area at bottom right?
        c_left = self.SLIDE_WIDTH - margin_x - Inches(2)
        c_top = self.SLIDE_HEIGHT - margin_y - Inches(2)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, c_left, c_top, Inches(4), Inches(4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.ACCENT_COLOR
        # Transparency hack via OXML later if needed. For now 100% solid might be too strong.
        # Let's skip the big circle to avoid visual clutter if transparency isn't easy.
        # Actually, let's do soft transparency hack.
        fill = circle.fill
        fill.solid()
        fill.fore_color.rgb = self.ACCENT_COLOR
        # Set alpha to ~10%
        # solidFill = fill.fore_color._xFill.solidFill
        # srgbClr = solidFill.srgbClr
        # srgbClr.append(OxmlElement('a:alpha', val="10000")) # 10%
        # This crashes sometimes depending on version. Simpler: skip for now or use light gray.
        slide.shapes._spTree.remove(circle._element) # Remove it, clean design preferred.

    def _add_page_num(self, slide, num, total):
        # Bottom right corner - Page number (matching reference image)
        # Symmetric margin: 0.5 inches from edge
        margin = Inches(0.5)
        box_width = Inches(1.0)
        num_left = self.SLIDE_WIDTH - margin - box_width
        num_bottom = self.SLIDE_HEIGHT - Inches(0.6)
        
        numBox = slide.shapes.add_textbox(num_left, num_bottom, box_width, Inches(0.4))
        p_num = numBox.text_frame.paragraphs[0]
        p_num.alignment = PP_ALIGN.RIGHT
        
        # Page number: 01 / 16
        run_num = p_num.add_run()
        run_num.text = f"{num:02d} / {total:02d}"
        run_num.font.name = self.FONT_MAIN
        run_num.font.size = Pt(16)
        run_num.font.color.rgb = self.TEXT_LIGHT
        run_num.font.bold = False

    def _add_logo(self, slide):
        """Adds the circular logo to the top left."""
        # renderer.py is in 'src', assets is in root.
        root_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        logo_path = os.path.join(root_dir, "assets", "logo_circle.png")
        
        if os.path.exists(logo_path):
            # Bottom-left corner positioning with symmetric margin (0.5 inches)
            size = Inches(0.8)
            margin = Inches(0.5)
            left = margin
            bottom = self.SLIDE_HEIGHT - Inches(0.9)
            pic = slide.shapes.add_picture(logo_path, left, bottom, width=size)
            # Add hyperlink
            pic.click_action.target_full_uri = "https://www.jxgqc.online"
        else:
            # For debug: we can print to console if file missing
            # print(f"DEBUG: Logo not found at {logo_path}")
            pass

    def create_title_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Blank
        self._set_bg(slide)
        self._add_card_container(slide)
        self._add_logo(slide)
        
        # Centering helper
        cx = self.SLIDE_WIDTH / 2
        cy = self.SLIDE_HEIGHT / 2
        
        # Title
        t_w = Inches(10)
        t_h = Inches(2)
        title_box = slide.shapes.add_textbox(cx - t_w/2, cy - Inches(1.5), t_w, t_h)
        p = title_box.text_frame.paragraphs[0]
        p.text = f"{self.subject}试题深度解析"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.FONT_MAIN
        p.font.size = Pt(60) # Scaled for PPT
        p.font.bold = True
        p.font.color.rgb = self.ACCENT_DARK
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(cx - t_w/2, cy + Inches(0.2), t_w, Inches(1))
        p = sub_box.text_frame.paragraphs[0]
        p.text = "—— 全真模拟 · 考点攻坚 ——"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.FONT_MAIN
        p.font.size = Pt(32)
        p.font.color.rgb = self.TEXT_LIGHT
        
        # Hint capsule
        hint_w = Inches(6)
        hint_h = Inches(0.8)
        hint_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            cx - hint_w/2, cy + Inches(2), hint_w, hint_h
        )
        hint_box.fill.solid()
        hint_box.fill.fore_color.rgb = RGBColor(248, 250, 252) # Slate 50
        hint_box.line.color.rgb = self.BG_COLOR
        
        p = hint_box.text_frame.paragraphs[0]
        p.text = "点击幻灯片交互：1次显答案，2次显解析"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.FONT_MAIN
        p.font.size = Pt(20)
        p.font.color.rgb = self.TEXT_HINT

        # Add signature to title as well
        self._add_page_num(slide, 0, 0) # Use helper, page digits will handle 00/00 or we can customize
        # Actually _add_page_num adds "00 / 00". For title, let's just use it or add a custom one.
        # Let's adjust _add_page_num to skip numbers if num is 0.

    def add_question_slides(self, questions):
        total = len(questions)
        for idx, q in enumerate(questions, 1):
            q_num = q['number']
            q_text_masked = q['question']
            answer_char = q['real_answer'] if q['real_answer'] else "?"
            options = q['options']
            expl_text = q['explanation']
            
            # --- Generate 3 Steps ---
            for step in range(1, 4):
                slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Blank
                self._set_bg(slide)
                self._add_card_container(slide)
                self._add_logo(slide)
                self._add_page_num(slide, idx, total)
                
                # --- Content Area ---
                margin_left = Inches(1.2)
                current_y = Inches(1.0) # Start top (Stem)
                content_width = Inches(11)

                # --- PIL MEASUREMENT ENGINE ---
                from PIL import ImageFont, ImageDraw, Image
                
                def measure_text_exact(text, font_size_pt, width_inches, font_path=None):
                    """
                    Uses PIL to simulate text wrapping and calculate exact height.
                    """
                    if font_path is None:
                        # renderer.py is in 'src', assets is in root.
                        root_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
                        font_path = os.path.join(root_dir, "assets", "msyh.ttc")
                        
                        # Fallback for local dev if missing in assets
                        if not os.path.exists(font_path):
                            font_path = "C:/Windows/Fonts/msyh.ttc"

                    font_size_px = int(font_size_pt * 1.333)
                    max_width_px = int(width_inches * 96)
                    
                    try:
                        # Use TTC index 0 for standard MS YaHei
                        font = ImageFont.truetype(font_path, font_size_px, index=0)
                    except IOError:
                        # Fallback to default if font not found
                        font = ImageFont.load_default()
                    
                    # Manual Word Wrap Logic
                    # We need to split text into lines that fit max_width_px
                    lines = []
                    
                    # Split by explicit newlines first
                    paragraphs = text.split('\n')
                    
                    for para in paragraphs:
                        current_line = ""
                        for char in para:
                            test_line = current_line + char
                            w = font.getbbox(test_line)[2] # width (right - left)
                            if w <= max_width_px:
                                current_line = test_line
                            else:
                                if current_line: lines.append(current_line)
                                current_line = char
                        if current_line: lines.append(current_line)
                    
                    num_lines = len(lines)
                    if num_lines < 1: num_lines = 1
                    
                    # Calculate Height
                    # Line spacing in PPT: we set Pt(34) for options (font 24). ~1.4x
                    # For Stem (font 26), line spacing is auto (~1.2x).
                    # Let's use proportional line height based on font size.
                    
                    if font_size_pt == 26: # Stem
                        line_height_pt = 36 # approx 1.4x 26
                    elif font_size_pt == 24: # Options
                        line_height_pt = 34 # Explicit requested
                    elif font_size_pt == 20: # Analysis
                        line_height_pt = 28 # approx 1.4x 20
                    else:
                        line_height_pt = font_size_pt * 1.4
                        
                    return Pt(num_lines * line_height_pt)

                # 1. STEM (Question Text)
                stem_est_h = measure_text_exact(q_text_masked, 26, 11)
                # Padding buffer just in case
                stem_est_h += Pt(10) 
                
                stem_box = slide.shapes.add_textbox(margin_left, current_y, content_width, stem_est_h)
                stem_box.text_frame.word_wrap = True
                stem_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                p = stem_box.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                
                self._render_stem(p, q_text_masked, answer_char, step, q.get('question_rich', []))
                
                current_y += stem_est_h + Pt(8)
                
                # 2. OPTIONS
                # Threshold calculation
                is_full_width = any([len(o) > 12 for o in options])
                # col_w is in Emu (Inches returns Emu)
                col_w = content_width if is_full_width else int(content_width / 2)
                
                option_rows = []
                if is_full_width:
                    for opt in options: option_rows.append([opt])
                else:
                    for i in range(0, len(options), 2):
                        row_opts = [options[i]]
                        if i + 1 < len(options): row_opts.append(options[i+1])
                        option_rows.append(row_opts)
                
                for row_opts in option_rows:
                    max_h = Pt(0)
                    for opt in row_opts:
                        # Convert EMU back to inches float for measurement
                        w_in = col_w / 914400.0
                        h = measure_text_exact(opt, 24, w_in)
                        if h > max_h: max_h = h
                    
                    if max_h < Inches(0.5): max_h = Inches(0.5) 
                    
                    for c_idx, opt in enumerate(row_opts):
                        x = int(margin_left + (c_idx * col_w))
                        
                        opt_box = slide.shapes.add_textbox(x, int(current_y), int(col_w), int(max_h))
                        opt_box.text_frame.word_wrap = True
                        opt_box.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                        
                        p = opt_box.text_frame.paragraphs[0]
                        p.text = opt
                        p.font.name = self.FONT_MAIN
                        p.font.size = Pt(24)
                        
                        # Highlight Option if it matches answer in Step 2+
                        # Option format usually "A. content" or "A.content"
                        # We check if it starts with "A." or "A ."
                        match_answer = False
                        if step >= 2 and answer_char:
                            prefix = f"{answer_char}."
                            if opt.strip().startswith(prefix):
                                match_answer = True
                        
                        if match_answer:
                            p.font.color.rgb = self.RED_ANSWER
                            p.font.bold = True
                        else:
                            p.font.color.rgb = self.TEXT_LIGHT
                            
                        p.line_spacing = Pt(34)
                    
                    current_y += max_h + Pt(5)

                # 3. ANALYSIS (Explanation)
                if step >= 3:
                     # Only show if not overflowing absurdly
                    current_y += Pt(10)
                    
                    an_height = measure_text_exact(expl_text, 20, 10.8)
                    an_height += Pt(20) # Padding
                    
                    # Decor Bar
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        margin_left, current_y, Inches(0.1), an_height
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = self.ACCENT_COLOR
                    bar.line.fill.background()
                    
                    # BG
                    bg = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        margin_left + Inches(0.1), current_y, content_width - Inches(0.1), an_height
                    )
                    bg.fill.solid()
                    bg.fill.fore_color.rgb = self.ANALYSIS_BG
                    bg.line.fill.background()
                    
                    # Text
                    txBox = slide.shapes.add_textbox(
                         margin_left + Inches(0.2), current_y + Inches(0.1), 
                         content_width - Inches(0.4), an_height - Inches(0.2)
                    )
                    txBox.text_frame.word_wrap = True
                    txBox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    p = txBox.text_frame.paragraphs[0]
                    p.text = "【解析】 " + expl_text
                    p.font.name = self.FONT_MAIN
                    p.font.size = Pt(20)
                    p.font.color.rgb = self.ACCENT_DARK 

    def _render_stem(self, p, q_text_masked, answer_char, step, question_rich=None):
        """Renders the question stem with interactive inline answer, supporting rich text formatting."""
        p.font.name = self.FONT_MAIN
        p.font.size = Pt(26) # Requested 26px
        p.font.color.rgb = self.TEXT_STEM
        p.font.bold = True
        
        # Split logic for inline bracket
        parts = q_text_masked.split('（   ）')
        # Simple rebuild
        if len(parts) > 1:
            # Part 1
            run = p.add_run()
            run.text = parts[0]
            
            # Bracket area
            run_open = p.add_run()
            run_open.text = "（ "
            run_open.font.color.rgb = self.TEXT_STEM
            
            run_ans = p.add_run()
            if step == 1:
                run_ans.text = "   " # Space
            else:
                run_ans.text = f"{answer_char}" # Answer
                run_ans.font.color.rgb = self.RED_ANSWER
                run_ans.font.bold = True
            
            run_close = p.add_run()
            run_close.text = " ）"
            run_close.font.color.rgb = self.TEXT_STEM
            
            # Remainder (handle multiple brackets? usually one answer)
            # Just take the rest joined
            run_rest = p.add_run()
            run_rest.text = "（   ）".join(parts[1:])
            
        else:
            # Fallback if no brackets
            p.text = q_text_masked
            if step >= 2:
                run = p.add_run()
                run.text = f"  （ {answer_char} ）"
                run.font.color.rgb = self.RED_ANSWER

    def save(self):
        self.prs.save(self.output_file)
